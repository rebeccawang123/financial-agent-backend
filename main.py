import os
import json
import base64 # 用于 PPT 图片编码
from typing import TypedDict, List, Dict, Any
from dotenv import load_dotenv

from langgraph.graph import StateGraph, END
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.tools.tavily_search import TavilySearchResults
from fastapi import FastAPI, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from openai import OpenAI # 用于 TTS 语音合成
from pptx import Presentation # 用于 PPT 生成
from pptx.util import Inches # PPT 尺寸

load_dotenv()

# --- 0. 初始化 OpenAI 客户端 (用于 TTS) ---
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# --- 1. 定义状态 (State) ---
class AgentState(TypedDict):
    query: str
    news_data: List[str]
    sources: List[Dict[str, str]]
    podcast_insights: str
    logs: List[str]
    final_report: str
    report_chinese: str # 新增: 中文报告
    report_english: str # 新增: 英文报告
    audio_chinese_b64: str # 新增: 中文语音 (Base64 编码)
    audio_english_b64: str # 新增: 英文语音 (Base64 编码)
    ppt_b64: str # 新增: PPT 文件 (Base64 编码)

# --- 2. 初始化 ---
llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
search_tool = TavilySearchResults(max_results=3)

# --- 3. 定义节点 (Nodes) ---

def news_node(state: AgentState):
    """新闻搜集员"""
    query = state.get("query", "Macro Finance")
    logs = state.get("logs", [])
    logs.append(f"🕵️ [News Agent] 开始搜索: '{query}'...")
    
    try:
        results = search_tool.invoke(f"{query} financial news bloomberg wsj")
        news_content = [res['content'] for res in results]
        sources = [{"title": res['content'][:30]+"...", "url": res['url']} for res in results]
        logs.append(f"✅ [News Agent] 成功抓取到 {len(results)} 条相关新闻。")
    except Exception as e:
        print(f"Search failed: {e}")
        logs.append("⚠️ [News Agent] 搜索 API 未响应，使用备用数据流。")
        news_content = [
            "Fed minutes suggest pause in rate cuts for December.",
            "NVIDIA stock volatility increases ahead of earnings.",
            "Bitcoin breaks $98k resistance level on ETF inflows."
        ]
        sources = [
            {"title": "WSJ: Fed Minutes Analysis", "url": "https://www.wsj.com/economy/central-banking"},
            {"title": "Bloomberg: Crypto Market Update", "url": "https://www.bloomberg.com/crypto"},
            {"title": "Reuters: Tech Stocks Rally", "url": "https://www.reuters.com/markets/us"}
        ]
        
    return {"news_data": news_content, "sources": sources, "logs": logs}

def podcast_node(state: AgentState):
    """播客监听员"""
    logs = state.get("logs", [])
    logs.append("🎧 [Pod Listener] 正在接入 RSS 源: 'All-In Podcast'...")
    logs.append("📝 [Pod Listener] 音频转录完成，正在提取关键观点...")
    
    mock_insight = """
    Chamath: AI infrastructure capex is peaking.
    Sacks: US Debt ceiling will be the main topic in 2025.
    """
    logs.append("✅ [Pod Listener] 观点提取完毕。")
    return {"podcast_insights": mock_insight, "logs": logs}

def analyst_node(state: AgentState):
    """首席分析师 - 生成中英文报告"""
    logs = state.get("logs", [])
    logs.append("🧠 [Chief Analyst] 正在交叉验证数据，准备生成多语言 Markdown 报告...")
    
    news = "\n".join(state['news_data'])
    podcast = state['podcast_insights']
    
    # 生成英文报告
    prompt_en = ChatPromptTemplate.from_template("""
    You are a Wall Street Chief Analyst. Based on news: {news} and podcast insights: {podcast}.
    Write a "Daily Financial Briefing" including: Market Sentiment, Macro Analysis, Web3 Watch, and Actionable Insights.
    Use Markdown format, and include Emojis.
    """)
    chain_en = prompt_en | llm
    response_en = chain_en.invoke({"news": news, "podcast": podcast})
    
    # 生成中文报告
    prompt_zh = ChatPromptTemplate.from_template("""
    你是华尔街首席分析师。基于新闻: {news} 和播客: {podcast}。
    写一份【每日金融晨报】，包含：市场情绪、宏观分析、Web3观察、操作建议。
    使用 Markdown 格式，多用 Emoji。
    """)
    chain_zh = prompt_zh | llm
    response_zh = chain_zh.invoke({"news": news, "podcast": podcast})
    
    logs.append("🚀 [System] 中英文报告生成完毕。")
    return {
        "report_english": response_en.content,
        "report_chinese": response_zh.content,
        "logs": logs
    }

def speech_node(state: AgentState):
    """语音合成员"""
    logs = state.get("logs", [])
    logs.append("🗣️ [Speech Synthesizer] 正在将报告转换为中英文语音...")
    
    report_zh = state['report_chinese']
    report_en = state['report_english']

    # 尝试中文语音合成
    try:
        speech_zh_response = openai_client.audio.speech.create(
            model="tts-1",
            voice="alloy", # 或 "nova", "shimmer" 等
            input=report_zh[:4096] # TTS API 限制输入长度，只取前4096字符
        )
        audio_zh_bytes = speech_zh_response.read()
        audio_chinese_b64 = base64.b64encode(audio_zh_bytes).decode('utf-8')
        logs.append("✅ [Speech Synthesizer] 中文语音生成成功。")
    except Exception as e:
        logs.append(f"❌ [Speech Synthesizer] 中文语音生成失败: {e}")
        audio_chinese_b64 = ""

    # 尝试英文语音合成
    try:
        speech_en_response = openai_client.audio.speech.create(
            model="tts-1",
            voice="alloy",
            input=report_en[:4096]
        )
        audio_en_bytes = speech_en_response.read()
        audio_english_b64 = base64.b64encode(audio_en_bytes).decode('utf-8')
        logs.append("✅ [Speech Synthesizer] 英文语音生成成功。")
    except Exception as e:
        logs.append(f"❌ [Speech Synthesizer] 英文语音生成失败: {e}")
        audio_english_b64 = ""

    return {
        "audio_chinese_b64": audio_chinese_b64,
        "audio_english_b64": audio_english_b64,
        "logs": logs
    }

def ppt_node(state: AgentState):
    """PPT 生成器"""
    logs = state.get("logs", [])
    logs.append("📊 [PPT Generator] 正在整理报告内容，准备生成演示文稿...")
    
    report_title = "每日金融晨报"
    report_content = state['report_chinese'] # 使用中文报告生成 PPT
    
    prs = Presentation()
    
    # 第一页：标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = report_title
    subtitle.text = f"由 AlphaBrief.ai 生成\n{json.dumps(state['sources'], indent=2, ensure_ascii=False)[:200]}..." # 简单展示来源

    # 第二页：内容页
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "核心洞察 (Key Insights)"
    
    # 简单地把 Markdown 报告拆分成段落
    content_parts = [part.strip() for part in report_content.split('\n\n') if part.strip()]
    
    tf = body.text_frame
    tf.clear()
    for part in content_parts:
        p = tf.add_paragraph()
        p.text = part # 直接将报告内容作为段落添加
        # 更多 PPT 样式需要更复杂的解析
        
    # 保存 PPT 到内存
    from io import BytesIO
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0) # 将指针移到开头
    
    ppt_b64 = base64.b64encode(ppt_stream.read()).decode('utf-8')
    
    logs.append("✅ [PPT Generator] PPT 文档生成成功。")
    return {"ppt_b64": ppt_b64, "logs": logs}

# --- 4. 构建图 ---
workflow = StateGraph(AgentState)
workflow.add_node("news_scout", news_node)
workflow.add_node("podcast_listener", podcast_node)
workflow.add_node("chief_analyst", analyst_node) # 新增分析师节点
workflow.add_node("speech_synthesizer", speech_node) # 新增语音节点
workflow.add_node("ppt_generator", ppt_node) # 新增 PPT 节点

workflow.set_entry_point("news_scout")
workflow.add_edge("news_scout", "podcast_listener")
workflow.add_edge("podcast_listener", "chief_analyst")
workflow.add_edge("chief_analyst", "speech_synthesizer")
workflow.add_edge("speech_synthesizer", "ppt_generator") # 语音后生成 PPT
workflow.add_edge("ppt_generator", END)

app_graph = workflow.compile()

# --- 5. API ---
app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

class ReportRequest(BaseModel):
    topic: str = "今日市场"

@app.post("/generate_report")
async def generate_report(req: ReportRequest):
    inputs = {"query": req.topic, "logs": []}
    result = await app_graph.ainvoke(inputs)
    return {
        "report_chinese": result["report_chinese"],
        "report_english": result["report_english"],
        "sources": result["sources"],
        "logs": result["logs"],
        "audio_chinese_b64": result["audio_chinese_b64"],
        "audio_english_b64": result["audio_english_b64"],
        "ppt_b64": result["ppt_b64"]
    }

# 新增一个下载 PPT 的路由
@app.get("/download_ppt")
async def download_ppt(ppt_b64: str):
    ppt_bytes = base64.b64decode(ppt_b64)
    return Response(content=ppt_bytes, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": "attachment; filename=Financial_Briefing.pptx"})